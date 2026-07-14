"""通用文档校验与轻量解析器。

PDF 和高级视觉引擎仍由 processor.py 负责；本模块负责 Office、Markdown、
HTML 的安全校验及结构化文本提取，避免上传路由继续把文档等同于 PDF。
"""

from __future__ import annotations

import os
import re
import zipfile
from pathlib import PurePosixPath


MAX_UPLOAD_BYTES = 100 * 1024 * 1024
MAX_ARCHIVE_MEMBERS = 10_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 300 * 1024 * 1024
MAX_ARCHIVE_COMPRESSION_RATIO = 250
MAX_TEXT_PREVIEW_CHARS = 300_000

DOCUMENT_TYPE_BY_EXTENSION = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".md": "markdown",
    ".markdown": "markdown",
    ".html": "html",
    ".htm": "html",
}

MIME_TYPE_BY_DOCUMENT_TYPE = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "markdown": "text/markdown",
    "html": "text/html",
}

ALLOWED_DECLARED_MIME_TYPES = {
    "pdf": {"application/pdf"},
    "docx": {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
        "application/zip",
        "application/x-zip-compressed",
    },
    "pptx": {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
        "application/zip",
        "application/x-zip-compressed",
    },
    "markdown": {"text/markdown", "text/plain", "text/x-markdown"},
    "html": {"text/html", "application/xhtml+xml", "text/plain"},
}

GENERIC_MIME_TYPES = {"", "application/octet-stream", "binary/octet-stream"}


class DocumentValidationError(ValueError):
    pass


def document_type_from_filename(filename: str) -> str:
    extension = os.path.splitext(filename or "")[1].lower()
    document_type = DOCUMENT_TYPE_BY_EXTENSION.get(extension)
    if not document_type:
        supported = ", ".join(DOCUMENT_TYPE_BY_EXTENSION)
        raise DocumentValidationError(f"不支持的文档格式；目前支持 {supported}")
    return document_type


def canonical_mime_type(document_type: str) -> str:
    return MIME_TYPE_BY_DOCUMENT_TYPE[document_type]


def decode_text_bytes(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        return data.decode("utf-16")
    if data.startswith(b"\xef\xbb\xbf"):
        return data.decode("utf-8-sig")
    for encoding in ("utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentValidationError("文本编码无法识别，请使用 UTF-8 或 GB18030")


def _validate_office_archive(file_path: str, document_type: str):
    required_member = "word/document.xml" if document_type == "docx" else "ppt/presentation.xml"
    try:
        with zipfile.ZipFile(file_path) as archive:
            members = archive.infolist()
            if len(members) > MAX_ARCHIVE_MEMBERS:
                raise DocumentValidationError("Office 文档包含过多压缩条目")

            member_names = {member.filename.replace("\\", "/") for member in members}
            if "[Content_Types].xml" not in member_names or required_member not in member_names:
                raise DocumentValidationError(f"文件内容不是有效的 {document_type.upper()} 文档")

            total_uncompressed = 0
            for member in members:
                normalized_name = member.filename.replace("\\", "/")
                member_path = PurePosixPath(normalized_name)
                if member_path.is_absolute() or ".." in member_path.parts:
                    raise DocumentValidationError("Office 压缩包包含不安全路径")
                if member.flag_bits & 0x1:
                    raise DocumentValidationError("暂不支持加密 Office 文档")

                total_uncompressed += member.file_size
                if total_uncompressed > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                    raise DocumentValidationError("Office 文档解压后体积过大")
                if member.file_size and member.compress_size == 0:
                    raise DocumentValidationError("Office 文档包含异常压缩条目")
                if member.compress_size:
                    ratio = member.file_size / member.compress_size
                    if ratio > MAX_ARCHIVE_COMPRESSION_RATIO:
                        raise DocumentValidationError("Office 文档压缩比异常，已拒绝处理")
    except zipfile.BadZipFile as exc:
        raise DocumentValidationError(f"文件内容不是有效的 {document_type.upper()} 文档") from exc


def validate_document_file(file_path: str, filename: str, declared_mime_type: str | None = None):
    document_type = document_type_from_filename(filename)
    file_size = os.path.getsize(file_path)
    if file_size <= 0:
        raise DocumentValidationError("上传文件为空")
    if file_size > MAX_UPLOAD_BYTES:
        raise DocumentValidationError("单个文件不能超过 100 MB")

    declared_mime = (declared_mime_type or "").split(";", 1)[0].strip().lower()
    if (
        declared_mime not in GENERIC_MIME_TYPES
        and declared_mime not in ALLOWED_DECLARED_MIME_TYPES[document_type]
    ):
        raise DocumentValidationError(
            f"文件扩展名与 MIME 类型不一致：{declared_mime or '未知'}"
        )

    if document_type == "pdf":
        with open(file_path, "rb") as file:
            if not file.read(5).startswith(b"%PDF-"):
                raise DocumentValidationError("文件内容不是有效的 PDF")
    elif document_type in {"docx", "pptx"}:
        _validate_office_archive(file_path, document_type)
    else:
        with open(file_path, "rb") as file:
            sample = file.read(min(file_size, 2 * 1024 * 1024))
        if b"\x00" in sample:
            raise DocumentValidationError("文本文件包含二进制内容")
        # 大文本抽样可能刚好截断多字节字符，去掉末尾最多四字节再做严格解码。
        decode_text_bytes(sample[:-4] if file_size > len(sample) and len(sample) > 4 else sample)

    return {
        "document_type": document_type,
        "mime_type": canonical_mime_type(document_type),
        "file_size": file_size,
    }


def _segment(text: str, section: str, location_type: str, location_index: int):
    clean_text = (text or "").strip()
    return {
        "text": clean_text,
        "section": (section or "正文").strip()[:180],
        "location_type": location_type,
        "location_index": max(0, int(location_index or 0)),
        "page": max(0, int(location_index or 0)) if location_type == "page" else 0,
    }


_MARKDOWN_HEADING_RE = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*#*\s*$")


def extract_markdown_segments_from_text(text: str, location_type: str = "heading"):
    segments = []
    current_section = "正文"
    current_lines = []
    next_index = 1

    def flush():
        nonlocal next_index
        content = "\n".join(current_lines).strip()
        if content:
            segments.append(_segment(content, current_section, location_type, next_index))
            next_index += 1

    for line in (text or "").splitlines():
        heading_match = _MARKDOWN_HEADING_RE.match(line)
        if heading_match:
            flush()
            current_lines = []
            current_section = heading_match.group(1).strip()
        current_lines.append(line)
    flush()

    if not segments and (text or "").strip():
        segments.append(_segment(text, "正文", location_type, 1))
    return segments


def extract_markdown_document(file_path: str):
    with open(file_path, "rb") as file:
        text = decode_text_bytes(file.read())
    segments = extract_markdown_segments_from_text(text)
    return {
        "segments": segments,
        "text": text,
        "unit_type": "heading",
        "unit_count": len(segments),
        "engine": "markdown",
    }


def _table_to_markdown(rows):
    normalized_rows = [
        [re.sub(r"\s+", " ", str(cell or "")).strip() for cell in row]
        for row in rows
    ]
    normalized_rows = [row for row in normalized_rows if any(row)]
    if not normalized_rows:
        return ""
    width = max(len(row) for row in normalized_rows)
    normalized_rows = [row + [""] * (width - len(row)) for row in normalized_rows]
    lines = ["| " + " | ".join(row) + " |" for row in normalized_rows]
    lines.insert(1, "| " + " | ".join(["---"] * width) + " |")
    return "\n".join(lines)


def extract_docx_document(file_path: str):
    try:
        from docx import Document
        from docx.document import Document as DocumentObject
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise RuntimeError("缺少 python-docx，请先安装 backend/requirements.txt") from exc

    document = Document(file_path)

    def iter_blocks(parent):
        if hasattr(parent, "iter_inner_content"):
            yield from parent.iter_inner_content()
            return
        parent_element = parent.element.body if isinstance(parent, DocumentObject) else parent._tc
        for child in parent_element.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, parent)
            elif child.tag == qn("w:tbl"):
                yield Table(child, parent)

    segments = []
    current_section = "正文"
    current_parts = []
    next_index = 1

    def flush():
        nonlocal next_index
        content = "\n\n".join(part for part in current_parts if part).strip()
        if content:
            segments.append(_segment(content, current_section, "heading", next_index))
            next_index += 1

    for block in iter_blocks(document):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style_name = (getattr(getattr(block, "style", None), "name", "") or "").lower()
            if style_name.startswith("heading") or style_name.startswith("标题"):
                flush()
                current_parts = []
                current_section = text
            current_parts.append(text)
        elif isinstance(block, Table):
            table_text = _table_to_markdown(
                [[cell.text for cell in row.cells] for row in block.rows]
            )
            if table_text:
                current_parts.append(table_text)
    flush()

    if not segments:
        raise RuntimeError("DOCX 中没有可提取的文本或表格")
    return {
        "segments": segments,
        "text": "\n\n".join(segment["text"] for segment in segments),
        "unit_type": "heading",
        "unit_count": len(segments),
        "engine": "python-docx",
    }


def extract_pptx_document(file_path: str):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，请先安装 backend/requirements.txt") from exc

    presentation = Presentation(file_path)
    segments = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = (slide.shapes.title.text or "").strip()
        parts = []
        for shape in sorted(slide.shapes, key=lambda item: (getattr(item, "top", 0), getattr(item, "left", 0))):
            if getattr(shape, "has_table", False):
                table_text = _table_to_markdown(
                    [[cell.text for cell in row.cells] for row in shape.table.rows]
                )
                if table_text:
                    parts.append(table_text)
            elif getattr(shape, "has_text_frame", False):
                shape_text = (shape.text or "").strip()
                if shape_text:
                    parts.append(shape_text)

        if getattr(slide, "has_notes_slide", False):
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            notes_text = (getattr(notes_frame, "text", "") or "").strip()
            if notes_text:
                parts.append(f"讲者备注：\n{notes_text}")

        text = "\n\n".join(dict.fromkeys(parts)).strip()
        if text:
            segments.append(_segment(text, title or f"幻灯片 {slide_index}", "slide", slide_index))

    if not segments:
        raise RuntimeError("PPTX 中没有可提取的文本、表格或备注")
    return {
        "segments": segments,
        "text": "\n\n".join(segment["text"] for segment in segments),
        "unit_type": "slide",
        "unit_count": len(presentation.slides),
        "engine": "python-pptx",
    }


def sanitize_html_snapshot(html_text: str):
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise RuntimeError("缺少 beautifulsoup4，请先安装 backend/requirements.txt") from exc

    soup = BeautifulSoup(html_text or "", "html.parser")
    for tag in soup.find_all(["script", "style", "iframe", "object", "embed", "form", "noscript"]):
        tag.decompose()
    for tag in soup.find_all(True):
        for attribute in list(tag.attrs):
            lowered = attribute.lower()
            if lowered.startswith("on") or lowered in {"srcdoc", "integrity"}:
                del tag.attrs[attribute]
        for attribute in ("href", "src", "action"):
            value = tag.attrs.get(attribute)
            if isinstance(value, str) and value.strip().lower().startswith(("javascript:", "data:")):
                del tag.attrs[attribute]
    return str(soup)


def extract_html_segments_from_text(html_text: str):
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise RuntimeError("缺少 beautifulsoup4，请先安装 backend/requirements.txt") from exc

    safe_html = sanitize_html_snapshot(html_text)
    soup = BeautifulSoup(safe_html, "html.parser")
    for tag in soup.find_all(["nav", "header", "footer", "aside"]):
        tag.decompose()
    root = soup.find("article") or soup.find("main") or soup.body or soup
    title = (soup.title.get_text(" ", strip=True) if soup.title else "") or "网页正文"
    current_section = title
    current_parts = []
    segments = []
    next_index = 1

    def flush():
        nonlocal next_index
        text = "\n\n".join(current_parts).strip()
        if text:
            segments.append(_segment(text, current_section, "web_section", next_index))
            next_index += 1

    for element in root.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "blockquote", "table"]):
        if element.name != "table" and element.find_parent("table") is not None:
            continue
        if element.name.startswith("h"):
            heading = element.get_text(" ", strip=True)
            if heading:
                flush()
                current_parts = []
                current_section = heading
            continue
        if element.name == "table":
            rows = [
                [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                for row in element.find_all("tr")
            ]
            text = _table_to_markdown(rows)
        else:
            text = element.get_text(" ", strip=True)
        if text and text not in current_parts:
            current_parts.append(text)
    flush()

    if not segments:
        fallback = root.get_text("\n", strip=True)
        if fallback:
            segments.append(_segment(fallback, title, "web_section", 1))
    return segments, safe_html, title


def extract_html_document(file_path: str):
    with open(file_path, "rb") as file:
        html_text = decode_text_bytes(file.read())
    segments, _, _ = extract_html_segments_from_text(html_text)
    if not segments:
        raise RuntimeError("HTML 中没有可提取的正文")
    return {
        "segments": segments,
        "text": "\n\n".join(segment["text"] for segment in segments),
        "unit_type": "web_section",
        "unit_count": len(segments),
        "engine": "html",
    }


def extract_native_document(file_path: str, document_type: str):
    if document_type == "docx":
        return extract_docx_document(file_path)
    if document_type == "pptx":
        return extract_pptx_document(file_path)
    if document_type == "markdown":
        return extract_markdown_document(file_path)
    if document_type == "html":
        return extract_html_document(file_path)
    raise ValueError(f"没有可用的原生解析器: {document_type}")


def preview_text(file_path: str, document_type: str):
    if document_type == "markdown":
        with open(file_path, "rb") as file:
            return decode_text_bytes(file.read())[:MAX_TEXT_PREVIEW_CHARS]
    if document_type == "html":
        result = extract_html_document(file_path)
        return result["text"][:MAX_TEXT_PREVIEW_CHARS]
    raise ValueError("该格式不支持文本预览")
